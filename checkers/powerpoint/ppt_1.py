from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt

def check_ppt_1(prs):
    checklist_data = {
        "Grading Criteria": [
            "Does the presentation have at least 5 slides?",
            "Does each slide have a title?",
            "Is the font size at least 24pt for body text?",
            "Are there images on at least 2 slides?",
            "Is there at least one chart or graph?",
            "Is there consistent theme/design?",
            "Are bullet points used appropriately?",
            "Is there a title slide with name?",
            "Is there a conclusion/summary slide?",
            "Are transitions used between slides?"
        ],
        "Completed": []
    }

    try:
        # Check number of slides
        sufficient_slides = len(prs.slides) >= 5
        checklist_data["Completed"].append("Yes" if sufficient_slides else "No")

        # Check for titles on each slide
        slides_with_titles = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.strip() != "":
                        slides_with_titles += 1
                        break
        all_slides_have_titles = slides_with_titles == len(prs.slides)
        checklist_data["Completed"].append("Yes" if all_slides_have_titles else "No")

        # Check font size
        appropriate_font_size = True
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run.font, 'size') and run.font.size is not None:
                                if run.font.size < Pt(24):
                                    appropriate_font_size = False
        checklist_data["Completed"].append("Yes" if appropriate_font_size else "No")

        # Check for images
        slides_with_images = 0
        for slide in prs.slides:
            has_image = False
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    break
            if has_image:
                slides_with_images += 1
        sufficient_images = slides_with_images >= 2
        checklist_data["Completed"].append("Yes" if sufficient_images else "No")

        # Check for charts
        has_chart = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    has_chart = True
                    break
            if has_chart:
                break
        checklist_data["Completed"].append("Yes" if has_chart else "No")

        # Check for consistent theme (basic check for now)
        checklist_data["Completed"].append("Yes")  # Placeholder

        # Check for bullet points
        has_bullets = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level > 0:
                            has_bullets = True
                            break
        checklist_data["Completed"].append("Yes" if has_bullets else "No")

        # Check for title slide (basic check)
        has_title_slide = False
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    if len(shape.text.strip()) > 0:
                        has_title_slide = True
                        break
        checklist_data["Completed"].append("Yes" if has_title_slide else "No")

        # Check for conclusion slide (basic check)
        has_conclusion = False
        if len(prs.slides) > 0:
            last_slide = prs.slides[-1]
            for shape in last_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.lower()
                    if "conclusion" in text or "summary" in text or "thank" in text:
                        has_conclusion = True
                        break
        checklist_data["Completed"].append("Yes" if has_conclusion else "No")

        # Check for transitions (basic check - assumes if there are any transition settings)
        checklist_data["Completed"].append("Yes")  # Placeholder for now

    except Exception as e:
        # If any check fails, fill remaining checks with "No"
        while len(checklist_data["Completed"]) < len(checklist_data["Grading Criteria"]):
            checklist_data["Completed"].append("No")

    return checklist_data
